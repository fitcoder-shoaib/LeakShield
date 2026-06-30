import re
import csv
import shutil
import subprocess
import tempfile
from email import policy
from email.parser import BytesParser
from pathlib import Path


PATTERNS = {
    "Credentials": {
        "Password": r"\b(?:password(?![_-]?hash)|pwd)\s*[:=\-]?\s*\S+",
        "Password Hash": r"\b(?:password_hash|passwd_hash|pwd_hash|bcrypt|argon2|sha(?:256|512)?)\s*[:=\-]?\s*\S+|\$2[aby]\$\d{2}\$[./A-Za-z0-9]{53}|\$argon2(?:id|i|d)\$[^\s]+|\b[a-f0-9]{64}\b|\b[a-f0-9]{128}\b",
        "PIN": r"\b(?:pin|atm_pin|security_pin)\s*[:=\-]?\s*\d{4,8}\b",
        "OTP": r"\b(?:otp|one[-_\s]?time[-_\s]?password|verification[-_\s]?code)\s*[:=\-]?\s*\d{4,8}\b",
        "API Key": r"\b(?:api[_-]?key|secret[_-]?key|client[_-]?secret)\s*[:=\-]?\s*[A-Za-z0-9_\-]{20,}",
        "Private Key": r"-----BEGIN (?:RSA |DSA |EC |OPENSSH |PGP )?PRIVATE KEY-----[\s\S]*?-----END (?:RSA |DSA |EC |OPENSSH |PGP )?PRIVATE KEY-----",
        "Access Token": r"\b(?:access[_-]?token|auth[_-]?token|bearer)\s*[:=\-]?\s*[A-Za-z0-9._\-]{20,}",
        "Refresh Token": r"\brefresh[_-]?token\s*[:=\-]?\s*[A-Za-z0-9._\-]{20,}",
    },
    "Identity": {
        "SSN": r"\b(?:(?:ssn|social[-_\s]?security[-_\s]?number)\s*[:#=\-]\s*)?(?!000|666|9\d{2})\d{3}-(?!00)\d{2}-(?!0000)\d{4}\b|\b(?:ssn|social[-_\s]?security[-_\s]?number)\s*[:#=\-]\s*(?!000|666|9\d{2})\d{3}\s?(?!00)\d{2}\s?(?!0000)\d{4}\b",
        "National ID": r"\b(?:national[_-]?id|national identification number|nid)\s*[:#=\-]\s*[A-Z0-9][A-Z0-9 \-]{5,24}\b",
        "Aadhaar Number": r"\b(?:aadhaar(?:_number)?|aadhar(?:_number)?|uidai)\s*[:#=\-]\s*[2-9]\d{3}[ -]?\d{4}[ -]?\d{4}\b",
        "Passport Number": r"\b(?:passport(?:_number)?|passport no\.?)\s*[:#=\-]?\s*[A-Z][0-9]{7}\b",
        "Driver License Number": r"\b(?:driver'?s?[_\s-]?license(?:_number)?|driving[_\s-]?licen[cs]e(?:_number)?|dl(?:_number)?)\s*[:#=\-]\s*[A-Z0-9][A-Z0-9 \-]{5,20}\b",
        "Tax ID": r"\b(?:tax[_\s-]?id|tin|ein|pan)\s*[:#=\-]?\s*(?:[A-Z]{5}[0-9]{4}[A-Z]|\d{2}-\d{7}|[A-Z0-9][A-Z0-9 \-]{6,20})\b",
    },
    "Financial": {
        "Bank Account Number": r"\b(?:bank[_\s-]?account(?:_number)?|account(?:_number)?|acct(?:_number)?)\s*[:#=\-]\s*\d{8,18}\b",
        "Credit Card Number": r"\b(?:credit[_\s-]?card(?:_number)?|card(?:_number)?|cc(?:_number)?)\s*[:#=\-]\s*(?:\d[ -]*?){13,19}\b|\b(?:\d{4}[- ]){3}\d{4}\b",
        "CVV": r"\b(?:cvv|cvc|card[_\s-]?security[_\s-]?code)\s*[:#=\-]\s*\d{3,4}\b",
        "Salary": r"\b(?:salary|annual[_\s-]?salary|monthly[_\s-]?salary|compensation|ctc)\s*[:#=\-]?\s*(?:[$₹€£]\s*)?\d[\d,]*(?:\.\d{2})?\b",
    },
    "Personal": {
        "Email": r"\b[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-z]{2,}\b",
        "Phone": r"\b(?:phone|mobile|cell|contact(?:_number)?)\s*[:#=\-]\s*(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{3,5}\)?[-.\s]?)?\d{3,5}[-.\s]?\d{4}\b|(?<!\w)\+91[-\s]?[6-9]\d{9}\b|\b[6-9]\d{9}\b|\b\(?\d{3}\)?[-.\s]\d{3}[-.\s]\d{4}\b",
        "Address": r"\b(?:address|home[_\s-]?address|billing[_\s-]?address|shipping[_\s-]?address)\s*[:#=\-]\s*[A-Za-z0-9][A-Za-z0-9,.\-/# ]{10,120}\b",
        "DOB": r"\b(?:dob|date[_\s-]?of[_\s-]?birth|birth[_\s-]?date)\s*[:#=\-]?\s*(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|\d{4}-\d{2}-\d{2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Sept|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4})\b",
    },
    "Health": {
        "Medical Record": r"\b(?:medical[_\s-]?record(?:_number)?|mrn|patient[_\s-]?id|health[_\s-]?record)\s*[:#=\-]\s*[A-Z0-9][A-Z0-9 \-]{4,24}\b",
        "Biometric Data": r"\b(?:fingerprint|finger[_\s-]?print|face[_\s-]?id|facial[_\s-]?recognition|retina(?:l)?[_\s-]?scan|iris[_\s-]?scan|voice[_\s-]?print|biometric(?:_data)?)\s*[:#=\-]?\s*\S+",
    },
    "Location": {
        "GPS Location": r"\b(?:gps(?:_location)?|latitude|longitude|lat|lng|lon)\s*[:#=\-]?\s*-?\d{1,3}\.\d{3,}\s*,?\s*(?:longitude|lng|lon)?\s*[:#=\-]?\s*-?\d{1,3}\.\d{3,}\b",
    },
}

WEIGHTS = {
    "Credentials": 6,
    "Identity": 5,
    "Financial": 5,
    "Personal": 3,
    "Health": 5,
    "Location": 3,
}
TEXT_EXTENSIONS = {".txt", ".log", ".json", ".csv"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}
SUPPORTED_EXTENSIONS = {
    *TEXT_EXTENSIONS,
    *IMAGE_EXTENSIONS,
    ".pdf",
    ".doc",
    ".docx",
    ".xlsx",
    ".pptx",
    ".eml",
    ".msg",
}
REDACTION_TOKEN = "[REDACTED:{label}]"


def analyze_document(text):
    findings, reasons = {}, []
    score = 0

    for category, items in PATTERNS.items():
        findings[category] = {}
        count = 0

        for label, pattern in items.items():
            matches = re.findall(pattern, text, flags=re.IGNORECASE)
            if matches:
                findings[category][label] = matches
                count += len(matches)

        if count:
            score += count * WEIGHTS[category]
            reasons.append(f"{count} {category} data item(s) detected")

    score = min(score, 100)
    level = "Low" if score == 0 else "Medium" if score <= 50 else "High"
    return findings, reasons, score, level


def generate_recommendations(findings):
    recs = []

    if findings.get("Credentials"):
        recs.append("Rotate exposed passwords or API keys immediately.")
    if findings.get("Financial"):
        recs.append("Mask or redact financial identifiers before sharing.")
    if findings.get("Identity"):
        recs.append("Treat identity documents as regulated data and verify disclosure approvals.")
    if findings.get("Personal"):
        recs.append("Anonymize personal data such as emails, phone numbers, addresses, and dates of birth.")
    if findings.get("Health"):
        recs.append("Handle medical and biometric data with strict access controls and audit logging.")
    if findings.get("Location"):
        recs.append("Remove precise location data unless it is required for the recipient.")
    if not recs:
        recs.append("No sensitive data detected. Document appears safe.")

    recs.append("Encrypt sensitive files and restrict access permissions.")
    return recs


def redact_text(text):
    for items in PATTERNS.values():
        for label, pattern in items.items():
            text = re.sub(
                pattern,
                REDACTION_TOKEN.format(label=label.replace(" ", "_").upper()),
                text,
                flags=re.IGNORECASE,
            )
    return text


def read_csv_text(file_or_path):
    if isinstance(file_or_path, (str, Path)):
        with open(file_or_path, newline="", encoding="utf-8", errors="ignore") as csv_file:
            rows = list(csv.reader(csv_file))
    else:
        raw_text = file_or_path.read().decode(errors="ignore")
        rows = list(csv.reader(raw_text.splitlines()))

    return "\n".join(" ".join(cell for cell in row if cell) for row in rows)


def read_xlsx_text(file_or_path):
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Install openpyxl to scan .xlsx files.") from exc

    workbook = load_workbook(file_or_path, data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(sheet.title)
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(" ".join(values))
    workbook.close()
    return "\n".join(parts)


def read_pptx_text(file_or_path):
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Install python-pptx to scan .pptx files.") from exc

    presentation = Presentation(file_or_path)
    parts = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def read_doc_text(path):
    path = Path(path)

    if shutil.which("textutil"):
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            check=True,
            capture_output=True,
            text=True,
        )
        return result.stdout

    if shutil.which("antiword"):
        result = subprocess.run(
            ["antiword", str(path)],
            check=True,
            capture_output=True,
            text=True,
        )
        return result.stdout

    raise RuntimeError("Install antiword, or use macOS textutil, to scan legacy .doc files.")


def read_eml_text(file_or_path):
    if isinstance(file_or_path, (str, Path)):
        with open(file_or_path, "rb") as eml_file:
            message = BytesParser(policy=policy.default).parse(eml_file)
    else:
        message = BytesParser(policy=policy.default).parsebytes(file_or_path.read())

    parts = [
        message.get("from", ""),
        message.get("to", ""),
        message.get("cc", ""),
        message.get("bcc", ""),
        message.get("subject", ""),
    ]

    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                parts.append(part.get_content())
    else:
        parts.append(message.get_content())

    return "\n".join(part for part in parts if part)


def read_msg_text(path):
    try:
        import extract_msg
    except ImportError as exc:
        raise RuntimeError("Install extract-msg to scan .msg files.") from exc

    message = extract_msg.Message(str(path))
    return "\n".join(
        part
        for part in [
            message.sender,
            message.to,
            message.cc,
            message.subject,
            message.body,
        ]
        if part
    )


def read_image_text(file_or_path):
    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError("Install pytesseract and pillow to OCR image files.") from exc

    try:
        return pytesseract.image_to_string(Image.open(file_or_path))
    except pytesseract.TesseractNotFoundError as exc:
        raise RuntimeError("Install the Tesseract OCR system binary to scan image files.") from exc


def read_path(path):
    path = Path(path)
    suffix = path.suffix.lower()

    if suffix in {".txt", ".log", ".json"}:
        return path.read_text(errors="ignore")

    if suffix == ".csv":
        return read_csv_text(path)

    if suffix == ".pdf":
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    if suffix == ".doc":
        return read_doc_text(path)

    if suffix == ".docx":
        from docx import Document

        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if suffix == ".xlsx":
        return read_xlsx_text(path)

    if suffix == ".pptx":
        return read_pptx_text(path)

    if suffix == ".eml":
        return read_eml_text(path)

    if suffix == ".msg":
        return read_msg_text(path)

    if suffix in IMAGE_EXTENSIONS:
        return read_image_text(path)

    raise ValueError(f"Unsupported file type: {path.suffix or 'no extension'}")


def read_uploaded_file(file):
    suffix = Path(file.name).suffix.lower()

    if suffix in {".txt", ".log", ".json"}:
        return file.read().decode(errors="ignore")

    if suffix == ".csv":
        return read_csv_text(file)

    if suffix == ".pdf":
        import pdfplumber

        with pdfplumber.open(file) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    if suffix == ".doc":
        with tempfile.NamedTemporaryFile(suffix=suffix) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            return read_doc_text(temp_file.name)

    if suffix == ".docx":
        from docx import Document

        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    if suffix == ".xlsx":
        return read_xlsx_text(file)

    if suffix == ".pptx":
        return read_pptx_text(file)

    if suffix == ".eml":
        return read_eml_text(file)

    if suffix == ".msg":
        with tempfile.NamedTemporaryFile(suffix=suffix) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            return read_msg_text(temp_file.name)

    if suffix in IMAGE_EXTENSIONS:
        return read_image_text(file)

    return ""


def scan_text(text):
    findings, reasons, score, level = analyze_document(text)
    return {
        "findings": findings,
        "reasons": reasons,
        "score": score,
        "risk": level,
        "recommendations": generate_recommendations(findings),
        "redacted": redact_text(text),
    }


def scan_path(path):
    path = Path(path)
    text = read_path(path)
    result = scan_text(text)
    result["file"] = str(path)
    return result
